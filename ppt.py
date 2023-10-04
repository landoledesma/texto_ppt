import pptx
from pptx.util import Inches, Pt
import os
import base64

TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def create_presentation(topic,slide_titles,slide_contents):
    if not os.path.exists('generated_ppt'):
        os.makedirs('generated_ppt')
    prs = pptx.Presentation()
    
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic 
    
    for slide_title,slide_content in zip(slide_titles,slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title 
        slide.shapes.placeholders[1].text = slide_content

        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE
    prs.save(f"generated_ppt/{topic}_presentation.pptx")

def download_link(topic):
    if not os.path.exists('generated_ppt'):
        os.makedirs('generated_ppt')
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'
