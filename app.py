import streamlit as st
import openai
import pptx
from  pptx.util import Inches, Pt
import base64
import os 
from dotenv import load_dotenv

# Cargando variables de entorno
load_dotenv("token.env")
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")
openai.api_key = os.getenv("OPENAI_API_KEY")

# format options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def slide_title(topic):
    prompt = f"Genera 5 titulos de slides para el tema dado '{topic}'."
    response = openai.Completion.create(
        engine = "text-davinci-003",
        prompt = prompt,
        max_tokens = 200
    )
    return response['choices'][0]['text'].split("\n")

def slide_content(slide_title):
    prompt = f"Genera contenido para eltitulo del slide '{slide_title}'."
    response = openai.Completion.create(
        engine = "text-davinci-003",
        prompt = prompt,
        max_tokens = 1000
    )
    return response['choices'][0]['text']

def create_presentation(topic,slide_titles,slide_contents):
    if not os.path.exists('generated_ppt'):
        os.makedirs('generated_ppt')
    prs = pptx.Presentation()
    
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic 
    
    for slide_title,slide_content in zip(slide_titles,slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title 
        slide.shapes.placeholders[1].text = slide_content

        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE
    prs.save(f"generated_ppt/{topic}_presentation.pptx")



def download_link(topic):
    # Verificar y crear el directorio 'generated_ppt' si no existe
    if not os.path.exists('generated_ppt'):
        os.makedirs('generated_ppt')
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"


    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


def main():
    st.title("PGeneracion de presentaciones ")

    topic = st.text_input("Escribe el tema de la presentacion:")
    generate_button = st.button("Crear presentacion")

    if generate_button and topic:
        st.info("Generando presentacion... Por favor espera .")
        slide_titles = slide_title(topic)
        filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
        print("Titulo: ", filtered_slide_titles)
        slide_contents = [slide_content(title) for title in filtered_slide_titles]
        print("Contenido: ", slide_contents)
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Todo salio correcto!")

        st.success("Presentacion generado correctamente!")
        st.markdown(download_link(topic), unsafe_allow_html=True)

if __name__ == "__main__":
    main()